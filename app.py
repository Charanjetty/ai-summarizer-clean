import streamlit as st
from summarizer import summarize_text
import docx2txt
import PyPDF2
import csv
import pandas as pd
import pptx
import io

# Page configuration
st.set_page_config(
    page_title="📄 AI Document Summarizer",
    page_icon="🧠",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom styling
st.markdown(
    """
    <style>
    .main {
        background-color: #f4f6f9;
        color: #333;
        font-family: 'Segoe UI', sans-serif;
    }
    .block-container {
        padding-top: 2rem;
        padding-bottom: 2rem;
    }
    .stButton>button {
        background-color: #4CAF50;
        color: white;
        font-weight: bold;
        border: none;
        padding: 10px 20px;
        border-radius: 8px;
    }
    .stButton>button:hover {
        background-color: #45a049;
    }
    .summary-box {
        background-color: #ffffff;
        color: #000000;  /* 👈 This line makes text visible */
        padding: 20px;
        border-radius: 12px;
        box-shadow: 0 4px 10px rgba(0, 0, 0, 0.1);
    }
    </style>
    """,
    unsafe_allow_html=True
)

# Header
st.title("📄 AI-Powered Document Summarizer")
st.markdown("Upload a document and click **Summarize** to generate an intelligent summary using Google's Gemini model.")

# Sidebar
with st.sidebar:
    st.header("💡 Instructions")
    st.markdown(
        """
        1. Upload a `.txt`, `.pdf`, `.docx`, `.pptx`, `.csv`, or `.xlsx` file.
        2. Click on **Summarize**.
        3. Get an instant summary powered by **Gemini 2.0 Flash**.
        4. Works best for large academic or technical documents.
        """
    )
    st.markdown("---")
    st.info("Multiple file formats supported now! 📂")

# File uploader
uploaded_file = st.file_uploader("📂 Upload your document", type=["txt", "pdf", "docx", "pptx", "csv", "xlsx"])

text = ""

if uploaded_file is not None:
    file_type = uploaded_file.name.split('.')[-1].lower()
    try:
        if file_type == "txt":
            text = uploaded_file.read().decode("utf-8")
        elif file_type == "pdf":
            reader = PyPDF2.PdfReader(uploaded_file)
            text = " ".join(page.extract_text() or "" for page in reader.pages)
        elif file_type == "docx":
            text = docx2txt.process(uploaded_file)
        elif file_type == "pptx":
            presentation = pptx.Presentation(uploaded_file)
            text = " ".join(shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text"))
        elif file_type == "csv":
            df = pd.read_csv(uploaded_file)
            text = df.to_string()
        elif file_type == "xlsx":
            df = pd.read_excel(uploaded_file)
            text = df.to_string()
        else:
            st.warning("Unsupported file type.")
    except Exception as e:
        st.error(f"❌ Error reading file: {e}")

if text:
    st.subheader("📃 Document Content")
    st.text_area("Your uploaded content:", text, height=300, label_visibility="collapsed")

    if st.button("✨ Summarize"):
        with st.spinner("Summarizing with Gemini..."):
            try:
                summary = summarize_text(text)
                st.subheader("🧠 Generated Summary")
                st.markdown(f"<div class='summary-box'>{summary}</div>", unsafe_allow_html=True)
            except Exception as e:
                st.error(f"❌ Error during summarization: {e}")
else:
    st.info("📁 Please upload a document to begin.")
    st.markdown("---")